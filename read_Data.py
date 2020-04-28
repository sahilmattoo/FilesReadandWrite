# -*- coding: utf-8 -*-
"""
Created on Tue Apr 28 15:01:34 2020

@author: smattoo5
"""

""" IMPORT LIBRARIES """

#!pip install slate3k
#!pip install docx2txt
import pandas as pd
import os
import time
import re
import docx2txt # docx parser
import slate3k as slate # pdf parser
from pptx import Presentation # ppt Parser
import shutil

## Ignore Logging issue reported due to pdf 
import logging
logging.propagate = False 
logging.getLogger().setLevel(logging.ERROR)


# Define DataFrame
Dataset = pd.DataFrame()

# Define Dictionary
# Define Varibale Data Dictionary having all the Documents
# Will Convert the Dictionary into Dataframe
DataDictionary ={}

#Set the path where all the folders with Docs are present
path = SET THE PATH FROM WHERE DOCUMENTS HAVE TO BE COPIED
movetopath = PATH WHERE YOU WANT THE FILES TO BE COPIED


def readWord(filename):
    #print(filename)
    text = docx2txt.process(filename)
    #print(my_text)
    return text
    #return docx2txt.process(filename)
    
#print(my_text)

def readPdf(filename):
    pdfText= ''
    with open(filename,'rb') as f:
        text = slate.PDF(f)
    for each in text:
        pdfText = pdfText+each
    return pdfText

def readPPTX(filename):
    pptxtext = ''
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pptxtext = pptxtext + shape.text
    return pptxtext

            
# Copy to Data Folder
def moveToData(path, movetopath):
    
    subfolderpaths = []
    
    # Save all the hierechy path details
    for folder in os.walk(path):
        subfolderpaths.append(folder[0])
    
    for eachpath in range(len(subfolderpaths)):
        folderpath = subfolderpaths[eachpath]
        #Find List of Directories present in Sub Folder
        entry = os.listdir(folderpath)
        for directory in entry:
            # If dicrectory is a file with extension
            if '.' in directory:
                #print('its a file')
                oldfilepath = folderpath+'\\'+directory
                #print(file)
                newfilepath = movetopath+'\\Data'
                shutil.copy(oldfilepath, newfilepath)
    
def parseInput(path):
    dataDictionary ={}
    
    path = movetopath+"//Data"
#    if os.path.exists(path):
#        os.rmdir(path)
#        
#    os.mkdir(path)
        
    try:
        
        for i in os.walk(path):
            counter = 0
            for j in i[2]:
                #print(j)
                counter  = counter+1
                doctype = j.split(".")[1]
                #print(doctype)
                if doctype == 'docx':
                    
                    filename = j
                    fn ="Data"+"\\"+filename
                    readDoc = readWord(fn)
                    dataDictionary.update({counter:[filename, readDoc, doctype]})
                    
                elif doctype == 'pdf':
                    filename = j
                    fn ="Data"+"\\"+filename
                    readDoc = readPdf(fn)
                    dataDictionary.update({counter:[filename, readDoc, doctype]})
                    
                elif doctype == 'pptx':
                    filename = j
                    fn ="Data"+"\\"+filename
                    readDoc = readPPTX(fn)
                    dataDictionary.update({counter:[filename, readDoc, doctype]})
                    
                ## List of Documents which we couldnot Parse    
                else:
                    filename = j
                    #print(filename)
                    #print("Cant Parse the File", filename, " for Document type as ", doctype)
                    oldfilepath = "Data"+"\\"+filename
                    #print(file)
                    newfilepath = movetopath+"\\NotRead"
                    shutil.copy(oldfilepath, newfilepath)
    
    except Exception as ex:
        print(ex)
        

        
    return dataDictionary



moveToData(path,movetopath)
DataDictionary = parseInput(movetopath) 


# Transform Dictionary to DataFrame
def Trans_to_dataframe(datadictionary):
    df = pd.DataFrame(datadictionary).transpose()
    df.columns=['Title', 'Description', 'DocType']
    return df

Dataset = Trans_to_dataframe(DataDictionary)
Dataset








